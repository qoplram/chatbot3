import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFaceHub
import speech_recognition as sr
from gtts import gTTS
import os
import requests
from bs4 import BeautifulSoup
from io import BytesIO
import pandas as pd
from docx import Document
from pptx import Presentation
import re  # Import the re module for regular expressions

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_text_from_url(url):
    try:
        # Send an HTTP GET request to the URL
        response = requests.get(url)

        # Check if the request was successful (status code 200)
        if response.status_code == 200:
            # Parse the HTML content of the page using BeautifulSoup
            soup = BeautifulSoup(response.content, 'html.parser')  # Use response.content
            # Extract text from the HTML (you may need to adjust this based on the webpage structure)
            # For simplicity, this example assumes all text within <p> tags is relevant.
            paragraphs = soup.find_all('p')
            text = ' '.join([p.get_text() for p in paragraphs])

            return text
        else:
            # Handle the case when the request was not successful
            return f"Failed to fetch content from {url}, Status Code: {response.status_code}"

    except Exception as e:
        # Handle exceptions that may occur during the request or parsing
        return f"An error occurred: {str(e)}"

def get_text_from_urls(urls):
    text = ""
    for url in urls:
        url_text = get_text_from_url(url)
        text += url_text
    return text

def get_text_chunks(text):
    if isinstance(text, bytes):
        # Decode the byte content assuming it's in UTF-8 encoding
        text = text.decode('utf-8')

    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
            
            # Convert bot responses to audio (only for the last message)
            if i == len(st.session_state.chat_history) - 1:
                response_text = message.content
                text_to_speech(response_text)

def speech_to_text():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        st.text("Listening...")
        r.adjust_for_ambient_noise(source)
        audio = r.listen(source)
        st.text("Processing...")

    try:
        user_question = r.recognize_google(audio)
        st.text(f"You said: {user_question}")
        handle_userinput(user_question)
    except sr.UnknownValueError:
        st.text("Sorry, I didn't understand your audio.")
    except sr.RequestError as e:
        st.text(f"Sorry, there was an error with the speech recognition service: {e}")

def text_to_speech(response_text):
    temp_audio_file = "temp_audio.mp3"
    tts = gTTS(text=response_text, lang='en')
    tts.save(temp_audio_file)
    st.audio(temp_audio_file, format='audio/mp3')

def process_uploaded_files(files, text_chunks):
    for file in files:
        if file:
            file_extension = file.name.split('.')[-1].lower()

            if file_extension == 'pdf':
                raw_text = get_pdf_text([file])
            elif file_extension == 'txt':
                raw_text = file.read()
            elif file_extension == 'docx':
                doc = Document(file)
                raw_text = "\n".join([p.text for p in doc.paragraphs])
            elif file_extension == 'pptx':
                prs = Presentation(file)
                raw_text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            raw_text += shape.text + "\n"
            elif file_extension == 'md':
                raw_text = file.read()
            elif file_extension == 'csv':
                df = pd.read_csv(file)
                raw_text = df.to_string(index=False)
            else:
                raw_text = ""

            text_chunks += get_text_chunks(raw_text)

def main():
    load_dotenv()
    st.set_page_config(page_title="Chat with multiple PDFs, URLs, and Files",
                       page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Chat with multiple PDFs, URLs, and Files :books:")

    st.sidebar.subheader("Upload Files")

    uploaded_files = st.sidebar.file_uploader("Upload your Files here", accept_multiple_files=True)

    if uploaded_files and st.sidebar.button("Start Processing Files"):
        with st.spinner("Processing Files..."):
            text_chunks = []
            process_uploaded_files(uploaded_files, text_chunks)
            vectorstore = get_vectorstore(text_chunks)
            st.session_state.conversation = get_conversation_chain(vectorstore)

    st.sidebar.subheader("Your URLs")
    url = st.sidebar.text_input("Enter a URL:")
    if st.sidebar.button("Add URL"):
        with st.spinner("Fetching content from URL..."):
            try:
                response = requests.get(url)
                response.raise_for_status()
                url_text = response.content.decode('utf-8')  # Decode the byte content
                text_chunks = get_text_chunks(url_text)
                vectorstore = get_vectorstore(text_chunks)
                st.session_state.conversation = get_conversation_chain(vectorstore)
            except Exception as e:
                st.error(f"Failed to fetch content from URL: {str(e)}")

    st.write("Ask a question about your documents or URLs:")

    user_question = st.text_input("", key="question_input")
    speak_button = st.button("Speak your question", key="speak_button", help="Click to ask a question by talking")

    if speak_button:
        speech_to_text()
    elif user_question:
        handle_userinput(user_question)

    if st.session_state.chat_history:
        response_text = st.session_state.chat_history[-1].content
        text_to_speech(response_text)

if __name__ == '__main__':
    main()